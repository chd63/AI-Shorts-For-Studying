import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
import video_script_runner
import sys
import json
import re
import os
import argparse
import pdfplumber
import random

# Load environment variables
load_dotenv()
api_key_1 = os.getenv("API_KEY")
if not api_key_1:
    raise ValueError("API_KEY not found. Check your .env file.")

# Set your API key
genai.configure(api_key=api_key_1)

# Function to read notes from a file
def read_notes_file(file_path):
    """
    Reads the content of a notes file and returns it as a string.
    """
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
        sys.exit(1)  # Exit the program if the file is missing

# Extract text from PowerPoint file
def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            content.append("\n".join(slide_text))
        return "\n\n".join(content)
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")
        sys.exit(1)

def extract_text_from_pdf(file_path):
    """
    Extracts text from a PDF file and returns it as a string.
    """
    try:
        content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                content.append(page.extract_text())
        return "\n\n".join(content)
    except Exception as e:
        print(f"Error processing PDF file: {e}")
        sys.exit(1)

# Function to get a random video file from the background_videos folder
def get_random_video():
    """Returns a random video file from the background_videos directory."""
    background_videos_dir = "./background_videos/"
    if not os.path.exists(background_videos_dir):
        raise FileNotFoundError(f"Error: The directory '{background_videos_dir}' does not exist.")

    video_files = [f for f in os.listdir(background_videos_dir) if f.endswith((".mp4", ".mov", ".avi", ".mkv"))]

    if not video_files:
        raise FileNotFoundError("Error: No video files found in the 'background_videos' directory.")

    return os.path.join(background_videos_dir, random.choice(video_files))

# Function to generate video based on input parameters
def generate_video(notes_text, words_per_screen, key_point, script_to_use, video_font_size, font_path, output_dir):
    """
    Generate videos based on the extracted text and key points.
    For each video generated, a different random background video is used.
    """
    # Request to extract key ideas from AI
    request_string = f"""
    Please extract {key_point} key ideas from the following text. 
    Each idea should be concise, around 100 words. 
    Format the output as a **valid Python list of strings**.

    Return only the JSON output without any additional text.

    Text:
    {notes_text}
    """

    # Initialize the model
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Generate text
    response = model.generate_content(request_string)

    response_text = response.text.strip()
    # Clean the response text
    response_text = re.sub(r"^```(?:json)?\n", "", response_text)
    response_text = re.sub(r"\n```$", "", response_text)

    # Safely convert the cleaned response into a Python list
    try:
        extracted_list = json.loads(response_text)  # Convert cleaned response into Python list
        if not isinstance(extracted_list, list):
            raise ValueError("Response is not a valid JSON list.")
    except (json.JSONDecodeError, ValueError) as e:
        print(f"❌ Error parsing response into list: {e}")
        extracted_list = [response_text]  # Fallback to raw response if parsing fails

    # Loop through the extracted key points and create a video for each
    for idx, key_point in enumerate(extracted_list, start=1):
        output_path = f"{output_dir}output_{idx}.mp4"  # Generate a unique output path for each key point
        print(f"\n🎬 Creating video for key point {idx}: {key_point}")

        # Pick a random video for each key point
        video_path = get_random_video()

        # Call the video creation function with each key point
        video_script_runner.run_video_script(
            script_to_use, video_path, key_point, output_path, font_path, video_font_size, words_per_screen
        )
        print(f"✅ Video {output_path} created successfully.")

# Modify the main function to handle PDF files as well
def main():
    parser = argparse.ArgumentParser(description="Process notes from a text file, PowerPoint, or PDF and generate videos.")
    parser.add_argument("notes_file", help="Path to the notes file (TXT, PPTX, or PDF)")
    parser.add_argument("key_point", type=int, help="Number of key points to extract and generate videos")
    parser.add_argument("words_per_screen", type=int, help="Number of words to display per screen in the video")
    parser.add_argument("-p", "--powerpoint", action="store_true", help="Use PowerPoint file instead of a text file")
    parser.add_argument("-d", "--pdf", action="store_true", help="Use PDF file instead of a text file")

    args = parser.parse_args()

    if args.powerpoint:
        notes_text = extract_text_from_pptx(args.notes_file)
    elif args.pdf:
        notes_text = extract_text_from_pdf(args.notes_file)  # PDF handling
    else:
        notes_text = read_notes_file(args.notes_file)

    # Set parameters (these can be adjusted by the user)
    script_to_use = "script"
    video_font_size = 20
    font_path = "/home/dev/.fonts/DejaVuSans.ttf"
    output_dir = "./generated_videos/"

    # Generate the video
    generate_video(notes_text, args.words_per_screen, args.key_point, script_to_use, video_font_size, font_path, output_dir)


if __name__ == "__main__":
    main()
